#!/usr/bin/env python3
"""Extrahiert Text aus PowerPoint-Dateien (.pptx) mit python-pptx.

Aufruf:
    python extract_pptx.py <pptx-datei-oder-ordner> [ausgabeordner]

- Ohne Ausgabeordner: Text auf stdout.
- Mit Ausgabeordner: pro .pptx eine .txt-Datei (UTF-8).
- Text enthaelt "--- Folie N ---"-Marker sowie Notizen (falls vorhanden).

Exit-Codes: 0 = ok, 1 = falscher Aufruf, 2 = nichts Nutzbares, 3 = python-pptx fehlt.
"""
import sys
import pathlib

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except AttributeError:
    pass

try:
    from pptx import Presentation
except ImportError:
    sys.stderr.write("PPTX_MISSING: 'python-pptx' ist nicht installiert. "
                     "Bitte 'pip install --user python-pptx' ausfuehren.\n")
    sys.exit(3)


def extract(path):
    prs = Presentation(str(path))
    parts = []
    chars = 0
    for num, slide in enumerate(prs.slides, start=1):
        parts.append("--- Folie %d ---" % num)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
                        chars += len(text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    line = " | ".join(cells)
                    parts.append(line)
                    chars += len(line.strip())
        # Notizen
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                parts.append("[Notizen] " + notes.strip())
                chars += len(notes.strip())
    return "\n".join(parts), len(prs.slides), chars


def main():
    if len(sys.argv) < 2:
        sys.stderr.write("usage: extract_pptx.py <pptx-or-folder> [out-dir]\n")
        sys.exit(1)

    src = pathlib.Path(sys.argv[1])
    out_dir = pathlib.Path(sys.argv[2]) if len(sys.argv) > 2 else None
    if out_dir is not None:
        out_dir.mkdir(parents=True, exist_ok=True)

    files = [src] if src.is_file() else sorted(src.glob("*.pptx"))
    if not files:
        sys.stderr.write("Keine .pptx-Dateien gefunden unter: %s\n" % src)
        sys.exit(0)

    usable = 0
    for f in files:
        try:
            content, nslides, chars = extract(f)
        except Exception as exc:
            sys.stderr.write("ERROR\t%s\t%s\n" % (f.name, exc))
            continue
        if out_dir is not None:
            target = out_dir / (f.stem + ".txt")
            target.write_text(content, encoding="utf-8")
            print("OK\t%s\t%d Folien, %d chars -> %s" % (f.name, nslides, chars, target.name))
        else:
            print(content)
        if chars > 0:
            usable += 1

    if usable == 0:
        sys.stderr.write("KEINE nutzbare Textquelle unter den PPTX.\n")
        sys.exit(2)


if __name__ == "__main__":
    main()
